import os
from fastapi import UploadFile, HTTPException

from PyPDF2 import PdfReader
import docx
from pptx import Presentation
from openpyxl import load_workbook
from PIL import Image
import pytesseract

UPLOAD_DIR = "uploads"

ALLOWED_EXTENSIONS = (
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".png",
    ".jpg",
    ".jpeg",
    ".txt",
)


def validate_file(file: UploadFile):
    if not file.filename.lower().endswith(ALLOWED_EXTENSIONS):
        raise HTTPException(status_code=400, detail="Unsupported file type")


async def save_file(file: UploadFile):
    os.makedirs(UPLOAD_DIR, exist_ok=True)
    file_path = os.path.join(UPLOAD_DIR, file.filename)

    with open(file_path, "wb") as f:
        f.write(await file.read())

    return file_path


# -------------------------
# TEXT EXTRACTION ENGINE
# -------------------------
def extract_text(file_path: str):

    ext = file_path.lower()

    # PDF
    if ext.endswith(".pdf"):
        reader = PdfReader(file_path)
        return " ".join(page.extract_text() or "" for page in reader.pages)

    # DOCX — extract paragraphs AND table cells (dates/experience often live in tables)
    elif ext.endswith(".docx"):
        doc = docx.Document(file_path)
        parts = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for p in cell.paragraphs:
                        if p.text.strip():
                            parts.append(p.text)
        return "\n".join(parts)

    # PPTX
    elif ext.endswith(".pptx"):
        prs = Presentation(file_path)
        text = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)

        return " ".join(text)

    # XLSX (Excel)
    elif ext.endswith(".xlsx"):
        wb = load_workbook(file_path)
        text = []

        for sheet in wb:
            for row in sheet.iter_rows(values_only=True):
                text.extend([str(cell) for cell in row if cell])

        return " ".join(text)

    # Images (OCR )
    elif ext.endswith((".png", ".jpg", ".jpeg")):
        image = Image.open(file_path)
        return pytesseract.image_to_string(image)

    # TXT
    elif ext.endswith(".txt"):
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    return ""
